import os
import re
import math
import PyPDF2
import textract
import speech_recognition as sr
import docx2txt as d2t
import tempfile
from pptx import Presentation
from langdetect import detect_langs
from summarizer import TransformerSummarizer
# from punctuator import Punctuator
from os import path
from pydub import AudioSegment
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError
)
from pdf2image import convert_from_path, convert_from_bytes


try:
    from PIL import Image
except ImportError:
    import Image
import pytesseract
pytesseract.tesseract_cmd = r'/opt/homebrew/Cellar/tesseract/5.2.0/bin/tesseract'


# def ocr_core(filename):
#     im = Image.open(filename)
#     im = im.point(lambda p: p > 75 and p + 100)
#     text = pytesseract.image_to_string(im)
#     lang = str(detect_langs(text)[0])[:2]
#     confidence = float(str(detect_langs(text)[0])[3:])
#     if lang == 'en':
#         while confidence < 0.99:
#             im = im.rotate(-90)
#             text = pytesseract.image_to_string(im)
#             confidence = float(str(detect_langs(text)[0])[3:])
#             print('Processing ...')
#     else:
#         print('No English Detected')
        
#     return text


# def punctuate(text):
#     p = Punctuator("C:\\Users\\lm44\\Documents\\Code\\Python\\Sumit Backend\\functions\\INTERSPEECH-T-BRNN.pcl")
#     punctuated = p.punctuate(text)
    
#     return punctuated


# def speech_to_text(filename):
    
#     r = sr.Recognizer()
#     with sr.AudioFile(filename) as source:
#         r.adjust_for_ambient_noise(source)
#         audio_text = r.record(source)
#         try:

#             text = r.recognize_google(audio_text)

#             print('Converting audio transcripts into text ...')
#             return(text)

#         except:
#              print('Sorry.. run again...')
                
                
# def mp3_to_wav(filename):
# #     get this function fixed
#     sound = AudioSegment.from_file(filename)
#     path = f'{filename[:-4]}.wav'
#     sound.export(path, format="wav")
    
#     return(path)


# # def pdf_to_text(filename):
# #     pdfFileObj = open(filename,'rb')
# #     pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
# #     num_pages = pdfReader.numPages
# #     count = 0
# #     text = ""

# #     while count < num_pages:
# #         pageObj = pdfReader.getPage(count)
# #         count +=1
# #         text += pageObj.extractText()

# #     if text != "":
# #         text = text
# #     else:
# #         text = textract.process(fileurl, method='tesseract', language='eng') 
    
# #     return text

# def pdf_to_text(filename):
#     images = convert_from_path(filename)
#     texts = []
#     for index, image in enumerate(images):
#         path = f'{filename[:-3]}image{index}.jpg'
#         image.save(path, 'JPEG')
#         text = ocr_core(path)
#         texts.append(text)
#         os.remove(path)

#     result = "".join(texts)
#     return result


# def worddoc_to_text(filename):
#     text = d2t.process(filename)
    
#     return text


# def ppt_to_text(filename):
#     texts = []
#     prs = Presentation(filename)
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, 'text'):
#                 texts.append(shape.text)
#     text = ''.join(texts)
    
#     return text


# # FILE = input("Enter Filepath - ")
# # FILE = FILE.replace("\'", "").replace("\"", "")
# # retention_ratio = float(input("Retained Information (eg. '0.3') - "))

# # if FILE[-4:] == '.png' or FILE[-4:] == '.jpg' or FILE[-5:] == '.jpeg'or FILE[-5:] == '.jfif':
# #     text = ocr_core(FILE)
    
# # elif FILE[-4:] == '.pdf':
# #     text = pdf_to_text(FILE)
    
# # elif FILE[-4:] == '.wav':
# #     text = speech_to_text(FILE)
# #     text = punctaute(text)
    
# # elif FILE[-4:] == '.mp3':
# #     wav = mp3_to_wav(FILE)
# #     text = speech_to_text(wav)
# #     text = punctuate(text)
    
# # elif FILE[-5:] == '.docx':
# #     text = worddoc_to_text(FILE) 
    
# # elif FILE[-5:] == '.pptx':
# #     text = ppt_to_text(FILE) 
    
# # elif FILE[-4:] == '.txt':
# #     with open(FILE, 'r') as file:
# #         text = file.read()

# # else:
# #     print('Error... Try Again')

# def sumit(notes, retention_ratio):
#     model = TransformerSummarizer(transformer_type="GPT2",transformer_model_key="gpt2-medium")
#     text = " ".join(notes.split())
#     summary = model(text, ratio=retention_ratio)
#     return " ".join(summary.split())


# # print("\nSummarizing ...\n")
# # print(punctuate(summary))

def ocr_core(filename):
    im = Image.open(filename)
    im = im.point(lambda p: p > 75 and p + 100)
    text = pytesseract.image_to_string(im)
    # lang = str(detect_langs(text)[0])[:2]
    # confidence = float(str(detect_langs(text)[0])[3:])
    # if lang == 'en':
    # while confidence < 0.99:
    # im = im.rotate(-90)
    text = pytesseract.image_to_string(im)
    # confidence = float(str(detect_langs(text)[0])[3:])
    print('Processing ...')
    # else:
    #     print('No English Detected')
    print(text)
    return text

ocr_core("/Users/liammoore/Documents/Code/Python/Sumit v2/Sumit Backend/functions/Screenshot 2023-02-11 at 23.01.19.png")

# def punctuate(text):
#     p = Punctuator("/Users/liammoore/Documents/Code/Python/Sumit v2/Sumit Backend/functions/INTERSPEECH-T-BRNN.pcl")
#     punctuated = p.punctuate(text)
    
#     return punctuated


def speech_to_text(filename):
    
    r = sr.Recognizer()
    with sr.AudioFile(filename) as source:
        r.adjust_for_ambient_noise(source)
        audio_text = r.record(source)
        try:

            text = r.recognize_google(audio_text)

            print('Converting audio transcripts into text ...')
            return(text)

        except:
             print('Sorry.. run again...')
                
                
def mp3_to_wav(filename):
#     get this function fixed
    sound = AudioSegment.from_file(filename)
    path = f'{filename[:-4]}.wav'
    sound.export(path, format="wav")
    
    return(path)


# def pdf_to_text(filename):
#     pdfFileObj = open(filename,'rb')
#     pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
#     num_pages = pdfReader.numPages
#     count = 0
#     text = ""

#     while count < num_pages:
#         pageObj = pdfReader.getPage(count)
#         count +=1
#         text += pageObj.extractText()

#     if text != "":
#         text = text
#     else:
#         text = textract.process(fileurl, method='tesseract', language='eng') 
    
#     return text

def pdf_to_text(filename):
    images = convert_from_path(filename)
    texts = []
    for index, image in enumerate(images):
        path = f'{filename[:-3]}image{index}.jpg'
        image.save(path, 'JPEG')
        text = ocr_core(path)
        texts.append(text)
        os.remove(path)

    result = "".join(texts)
    return result


def worddoc_to_text(filename):
    text = d2t.process(filename)
    
    return text


def ppt_to_text(filename):
    texts = []
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
    text = ''.join(texts)
    
    return text


# FILE = input("Enter Filepath - ")
# FILE = FILE.replace("\'", "").replace("\"", "")
# retention_ratio = float(input("Retained Information (eg. '0.3') - "))

# if FILE[-4:] == '.png' or FILE[-4:] == '.jpg' or FILE[-5:] == '.jpeg'or FILE[-5:] == '.jfif':
#     text = ocr_core(FILE)
    
# elif FILE[-4:] == '.pdf':
#     text = pdf_to_text(FILE)
    
# elif FILE[-4:] == '.wav':
#     text = speech_to_text(FILE)
#     text = punctaute(text)
    
# elif FILE[-4:] == '.mp3':
#     wav = mp3_to_wav(FILE)
#     text = speech_to_text(wav)
#     text = punctuate(text)
    
# elif FILE[-5:] == '.docx':
#     text = worddoc_to_text(FILE) 
    
# elif FILE[-5:] == '.pptx':
#     text = ppt_to_text(FILE) 
    
# elif FILE[-4:] == '.txt':
#     with open(FILE, 'r') as file:
#         text = file.read()

# else:
#     print('Error... Try Again')

def sumit(notes, retention_ratio):
    model = TransformerSummarizer(transformer_type="GPT2",transformer_model_key="gpt2-medium")
    text = " ".join(notes.split())
    summary = model(text, ratio=retention_ratio)
    return " ".join(summary.split())


# print("\nSummarizing ...\n")
# print(punctuate(summary))



